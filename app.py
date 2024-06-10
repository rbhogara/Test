import streamlit as st
import io
import PyPDF2
import docx
import pptx
from langchain_community.vectorstores import Chroma
from langchain_community import embeddings
from langchain_community.llms import Ollama
import time
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate
from langchain.text_splitter import CharacterTextSplitter
from collections import namedtuple
import pandas as pd
import matplotlib.pyplot as plt
import threading

def process_input(uploaded_files):
    # Extract text from uploaded files
    docs_list = []
    for uploaded_file in uploaded_files:
        with io.BytesIO(uploaded_file.read()) as f:
            try:
                if uploaded_file.name.endswith('.pdf'):
                    pdf_reader = PyPDF2.PdfReader(f)
                    text = ""
                    for page in range(len(pdf_reader.pages)):
                        text += pdf_reader.pages[page].extract_text()
                    docs_list.append(text)
                elif uploaded_file.name.endswith('.docx'):
                    doc = docx.Document(f)
                    text = ""
                    for para in doc.paragraphs:
                        text += para.text
                    docs_list.append(text)
                elif uploaded_file.name.endswith('.pptx'):
                    presentation = pptx.Presentation(f)
                    text = ""
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text += shape.text
                    docs_list.append(text)
                else:
                    st.error(f"Unsupported file type: {uploaded_file.name}")
                    continue
            except Exception as e:
                st.error(f"Error processing file: {uploaded_file.name}")
                st.error(str(e))
                continue

    # Split the text into chunks
    text_splitter = CharacterTextSplitter(
            separator="\n",
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )    
    chunks = text_splitter.split_text(" ".join(docs_list))

    
    # Convert text chunks into embeddings and store in vector database

    Document = namedtuple('Document', ['page_content', 'metadata'])

    # Create Document objects with page content and metadata
    documents = [Document(page_content=chunk, metadata={}) for chunk in chunks]

    # Pass the list of Document objects to Chroma.from_documents()
    vectorstore = Chroma.from_documents(
        documents=documents,
        collection_name="rag-chroma",
        embedding=embeddings.OllamaEmbeddings(model='nomic-embed-text'),
    )

    retriever = vectorstore.as_retriever()
    return retriever

def question(question, retriever):
    model_local = Ollama(model="mistral")

    # Perform the RAG
    after_rag_template = """Answer the question based only on the following context:
    {context}
    Question: {question}
    """
    after_rag_prompt = ChatPromptTemplate.from_template(after_rag_template)
    after_rag_chain = (
        {"context": retriever, "question": RunnablePassthrough()}
        | after_rag_prompt
        | model_local
        | StrOutputParser()
    )
    return after_rag_chain.invoke(question)

def process_csv(uploaded_file):
    try:
        df = pd.read_csv(uploaded_file)
        return df
    except Exception as e:
        st.error(f"Error processing CSV file: {uploaded_file.name}")
        st.error(str(e))
        return None
    
def plot_data(df):
    st.write("Select columns to plot:")
    columns = st.multiselect("Columns", df.columns)
    if st.button('Plot'):
        if len(columns) > 1:
            for column in columns:
                if df[column].dtype == 'object':
                    # Categorical data
                    fig, ax = plt.subplots()
                    ax.bar(df[column].value_counts().index, df[column].value_counts().values)
                    ax.set_xlabel('Category')
                    ax.set_ylabel('Count')
                    ax.set_title('Categorical Data')
                    st.pyplot(fig)
                elif df[column].dtype in ['int64', 'float64']:
                    # Numerical data
                    fig, ax = plt.subplots()
                    ax.hist(df[column], bins=50)
                    ax.set_xlabel('Value')
                    ax.set_ylabel('Frequency')
                    ax.set_title('Numerical Data')
                    st.pyplot(fig)
                else:
                    st.warning("Unsupported data type.")
        else:
            column = columns
            if df[column].dtype == 'object':
                # Categorical data
                fig, ax = plt.subplots()
                ax.bar(df[column].value_counts().index, df[column].value_counts().values)
                ax.set_xlabel('Category')
                ax.set_ylabel('Count')
                ax.set_title('Categorical Data')
                st.pyplot(fig)
            elif df[column].dtype in ['int64', 'float64']:
                # Numerical data
                fig, ax = plt.subplots()
                ax.hist(df[column], bins=50)
                ax.set_xlabel('Value')
                ax.set_ylabel('Frequency')
                ax.set_title('Numerical Data')
                st.pyplot(fig)
            else:
                st.warning("Unsupported data type.")



def main():
    st.title("DocuTalk")
    st.write("Upload files and enter a question to query the documents.")

    uploaded_files = st.file_uploader("Upload files", type=["pdf", "docx", "pptx", "csv"], accept_multiple_files=True)

    if uploaded_files:
        csv_files = [file for file in uploaded_files if file.name.endswith('.csv')]
        if csv_files:
            st.write("CSV file detected. Plotting data...")
            for file in csv_files:
                df = process_csv(file)
                if df is not None:
                    plot_data(df)
        else:
            if st.button('Generate Embeddings'):
                with st.spinner('Generating Embeddings...'):
                    retriever = process_input(uploaded_files)
                    st.success('Embeddings generated successfully!')

            question_input = st.text_input("Question")

            if st.button('Query Documents'):
                with st.spinner('Processing...'):
                    retriever = process_input(uploaded_files)
                    answer = question(question_input, retriever)
                    st.text_area("Answer", value=answer, height=300, disabled=True)
    else:
        st.warning("Please upload at least one file.")

if __name__ == "__main__":
    main()
